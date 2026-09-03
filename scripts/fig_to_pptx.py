"""Editable PPTX export, one file per paper figure (figs_pptx/<name>.pptx).

Diagram figures: the live matplotlib artists (rounded boxes, dashed groups, circles,
arrows, embedded photos, every text label) are walked and re-emitted as native
PowerPoint shapes at the same positions, so each element can be moved/edited.
Chart figures: rebuilt as native PPT charts with editable data sheets.
"""
import json, os, sys, tempfile
import numpy as np
import matplotlib; matplotlib.use("Agg")
import matplotlib.figure
from matplotlib.colors import to_hex
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from lxml import etree

R = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTD = os.path.join(R, "figs_pptx")
os.makedirs(OUTD, exist_ok=True)
SW, SH = 13.333, 7.5

DIAGRAMS = {"fig_wm_failure", "fig_operations", "fig_view_quality", "fig_paths_example",
            "fig_overview", "fig_action_space", "fig_gate_example", "fig_control_prompt",
            "fig_stopping", "fig_search_example", "fig_training_pipeline",
            "fig_oracle_example", "fig_runtime"}

def sanitize(t):
    for a, b in [("$\\approx$", "≈"), ("$^\\circ$", "°"), ("$\\times$", "×"),
                 ("$\\checkmark$", "✓"), ("$\\le$", "≤"), ("$\\rightarrow$", "→"),
                 ("\\,", " "), ("$", ""), ("\\%", "%"), ("--", "–"), ("``", '"'), ("''", '"')]:
        t = t.replace(a, b)
    return t

def rgb(c):
    try:
        h = to_hex(c)
    except Exception:
        return None
    return RGBColor.from_string(h.lstrip("#")[:6].upper())

def is_dashed(art):
    ls = art.get_linestyle()
    return not (ls in ("-", "solid") or ls is None)


def set_dash(line):
    ln = line._get_or_add_ln()
    d = ln.find("{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash")
    if d is None:
        d = etree.SubElement(ln, "{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash")
    d.set("val", "dash")

def add_arrowhead(line_elem):
    ln = line_elem
    tail = ln.find("{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd")
    if tail is None:
        tail = etree.SubElement(ln, "{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd")
    tail.set("type", "triangle"); tail.set("w", "med"); tail.set("len", "med")

def new_slide():
    prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])

def add_textbox(slide, x, y, text, pt, color, ha="center", va="top", bold=False, rot=0.0):
    va = {"baseline": "bottom", "center_baseline": "center"}.get(va, va)
    ha = ha if ha in ("left", "center", "right") else "center"
    text = sanitize(text)
    lines = text.split("\n")
    w = max(0.5, max(len(l) for l in lines) * pt * 0.0090)
    h = max(0.22, len(lines) * pt * 0.0195)
    left = {"center": x - w / 2, "left": x, "right": x - w}[ha]
    top = {"top": y, "center": y - h / 2, "bottom": y - h}[va]
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = False
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "center": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[va]
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, l in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = l
        p.alignment = {"center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT}[ha]
        for r in (p.runs or [p.add_run()]):
            r.font.size = Pt(max(6, pt)); r.font.bold = bold
            if color is not None: r.font.color.rgb = color
    if rot: tb.rotation = (360 - rot) % 360
    return tb

def convert_fig(fig, name):
    ax = fig.axes[0]
    x0, x1 = ax.get_xlim(); y0, y1 = ax.get_ylim()
    W, H = x1 - x0, y1 - y0
    s = min(12.6 / W, 6.9 / H)
    ox, oy = (SW - W * s) / 2, (SH - H * s) / 2
    X = lambda x: ox + (x - x0) * s
    Y = lambda y: oy + (y1 - y) * s
    prs, slide = new_slide()
    arts = [a for a in ax.get_children() if a is not ax.patch]
    arts.sort(key=lambda a: getattr(a, "get_zorder", lambda: 0)())
    tmpn = 0
    for art in arts:
        cls = type(art).__name__
        if cls == "AxesImage":
            arr = np.asarray(art.get_array())
            if arr.dtype != np.uint8:
                arr = (np.clip(arr, 0, 1) * 255).astype(np.uint8)
            l, r_, b, t = art.get_extent()
            tmp = os.path.join(tempfile.gettempdir(), f"p2p_{name}_{tmpn}.png"); tmpn += 1
            Image.fromarray(arr).save(tmp)
            slide.shapes.add_picture(tmp, Inches(X(l)), Inches(Y(t)),
                                     width=Inches(abs(r_ - l) * s), height=Inches(abs(t - b) * s))
        elif cls == "FancyBboxPatch":
            bx, by, bw, bh = art.get_x(), art.get_y(), art.get_width(), art.get_height()
            sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(X(bx)), Inches(Y(by + bh)),
                                        Inches(bw * s), Inches(bh * s))
            try: sp.adjustments[0] = 0.08
            except Exception: pass
            fc = art.get_facecolor()
            if fc is None or (len(fc) == 4 and fc[3] == 0):
                sp.fill.background()
            else:
                sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fc)
            ec = rgb(art.get_edgecolor())
            if ec is None: sp.line.fill.background()
            else:
                sp.line.color.rgb = ec
                sp.line.width = Pt(max(0.75, art.get_linewidth()))
                if is_dashed(art): set_dash(sp.line)
            sp.shadow.inherit = False
            sp.text_frame.word_wrap = True
        elif cls == "Circle":
            cx, cy = art.center; rr = art.radius
            sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(X(cx - rr)), Inches(Y(cy + rr)),
                                        Inches(2 * rr * s), Inches(2 * rr * s))
            sp.fill.solid(); sp.fill.fore_color.rgb = rgb(art.get_facecolor())
            ec = rgb(art.get_edgecolor())
            if ec is not None:
                sp.line.color.rgb = ec; sp.line.width = Pt(max(0.75, art.get_linewidth()))
            else:
                sp.line.fill.background()
            sp.shadow.inherit = False
        elif cls == "FancyArrowPatch":
            try: (pa, pb) = art._posA_posB
            except Exception: continue
            conn = slide.shapes.add_connector(2, Inches(X(pa[0])), Inches(Y(pa[1])),
                                              Inches(X(pb[0])), Inches(Y(pb[1])))
            col = rgb(art.get_edgecolor())
            if col is not None: conn.line.color.rgb = col
            conn.line.width = Pt(max(1.0, art.get_linewidth()))
            if is_dashed(art): set_dash(conn.line)
            ln = conn.line._get_or_add_ln()
            add_arrowhead(ln)
            conn.shadow.inherit = False
        elif cls == "Text":
            t = art.get_text()
            if not t.strip(): continue
            px, py = art.get_position()
            add_textbox(slide, X(px), Y(py), t, art.get_fontsize() * s,
                        rgb(art.get_color()), art.get_horizontalalignment(),
                        art.get_verticalalignment(), str(art.get_fontweight()) in ("bold", "700"),
                        float(art.get_rotation()))
    prs.save(os.path.join(OUTD, name + ".pptx"))
    print("diagram", name, "shapes:", len(slide.shapes))

# ---------------- native charts ----------------
INK, RED, TEALD, BLUE, GREY, AMB = "1F2328", "D62828", "155E63", "0969DA", "9AA4B1", "B26A00"

def style_series(plot_or_chart, idx, hexcol):
    ser = plot_or_chart.series[idx]
    ser.format.fill.solid(); ser.format.fill.fore_color.rgb = RGBColor.from_string(hexcol)

def chart_slide(name, builder):
    prs, slide = new_slide()
    builder(slide)
    prs.save(os.path.join(OUTD, name + ".pptx"))
    print("chart  ", name)

def col_chart(slide, x, y, w, h, cats, series, title=None, stacked=False, bar=False):
    cd = CategoryChartData(); cd.categories = cats
    for nm, vals in series: cd.add_series(nm, vals)
    ct = XL_CHART_TYPE.BAR_STACKED if (stacked and bar) else (
        XL_CHART_TYPE.COLUMN_STACKED if stacked else XL_CHART_TYPE.COLUMN_CLUSTERED)
    gf = slide.shapes.add_chart(ct, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.has_legend = len(series) > 1
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.TOP; ch.legend.include_in_layout = False
    if title:
        ch.has_title = True; ch.chart_title.text_frame.text = title
    else:
        ch.has_title = False
    return ch

def b_eff_methods(slide):
    cats = ["Direct input", "Static memory", "Tree d=1 (non-adpt.)", "ViewTree (ours)"]
    ch = col_chart(slide, 1.2, 0.9, 10.9, 5.6, cats,
                   [("latency (s)", (11.7, 11.4, 46.3, 24.5)), ("energy (kJ)", (0.529, 0.518, 2.110, 1.104))])
    style_series(ch.plots[0], 0, TEALD); style_series(ch.plots[0], 1, RED)
    add_textbox(slide, SW / 2, 6.7, "Per-question cost by method (Jetson AGX Orin, 7B; energy in kJ on the same axis)", 14, None)

def b_eff_routes(slide):
    cats = ["depth 0 (gated)  71%", "depth 1  15%", "depth 2  6%", "depth 3  8%"]
    ch = col_chart(slide, 1.2, 0.9, 10.9, 5.6, cats, [("latency (s)", (9.6, 26.0, 70.9, 117.2))])
    style_series(ch.plots[0], 0, TEALD)
    add_textbox(slide, SW / 2, 6.7, "Adaptive route mix - expected mix 24.5 s; one 16-frame call 11.7 s", 14, None)

def b_eff_breakdown(slide):
    cats = ["ViewTree depth 1", "ViewTree depth 0", "Direct (16 frames)"]
    series = [("gate", (4.33, 4.33, 0)), ("render+encode", (0.6, 0, 0)),
              ("answer call(s)", (16.8, 5.19, 11.7)), ("control call", (4.3, 0, 0))]
    ch = col_chart(slide, 1.2, 0.9, 10.9, 5.6, cats, series, stacked=True, bar=True)
    for i, c in enumerate([BLUE, AMB, TEALD, RED]): style_series(ch.plots[0], i, c)
    add_textbox(slide, SW / 2, 6.7, "Where the time goes (s per question)", 14, None)

def b_complexity(key, xlab):
    bins = json.load(open(os.path.join(R, "results/paperfill/complexity_bins.json")))["bins"][key]
    def b(slide):
        cats = [x["label"] for x in bins]
        ch = col_chart(slide, 1.2, 0.9, 10.9, 5.6, cats,
                       [("Direct input", tuple(x["direct"] for x in bins)),
                        ("Standard SFT", tuple(x["sft"] for x in bins)),
                        ("ViewTree", tuple(x["vt"] for x in bins))])
        for i, c in enumerate([GREY, TEALD, RED]): style_series(ch.plots[0], i, c)
        add_textbox(slide, SW / 2, 6.7, f"Score vs {xlab} (VSI-Bench held-out, 7B)", 14, None)
    return b

def b_frames_scaling(slide):
    fr = (16, 24, 32, 48, 64, 96, 128, 192, 256, 384)
    lat = (12.3, 16.6, 21.0, 31.2, 40.8, 61.6, 83.4, 129.9, 179.2, 299.7)
    mem = (22.6, 23.1, 23.5, 24.5, 25.4, 27.2, 29.1, 32.8, 36.5, 43.8)
    for k, (vals, ttl, col, xoff) in enumerate([
            (lat, "latency per question (s)", RED, 0.5), (mem, "peak GPU memory (GB)", TEALD, 6.9)]):
        cd = XyChartData(); srs = cd.add_series(ttl)
        for a, b in zip(fr, vals): srs.add_data_point(a, b)
        gf = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER_LINES, Inches(xoff), Inches(1.0),
                                    Inches(6.0), Inches(5.2), cd)
        ch = gf.chart; ch.has_legend = False; ch.has_title = True
        ch.chart_title.text_frame.text = ttl
        ser = ch.plots[0].series[0]
        ser.format.line.color.rgb = RGBColor.from_string(col); ser.format.line.width = Pt(2)
    add_textbox(slide, SW / 2, 6.6, "Frame-count scaling on Jetson (measured); OOM at 512 frames, position-embedding limit ~128", 13, None)

CHARTS = {
    "fig_frames_scaling": b_frames_scaling,
    "vt_eff_methods": b_eff_methods,
    "vt_eff_routes": b_eff_routes,
    "vt_eff_breakdown": b_eff_breakdown,
    "vt_object_number": b_complexity("object_number", "ground-truth object count"),
    "vt_spatial_scale": b_complexity("spatial_scale", "room area"),
    "vt_temporal_duration": b_complexity("temporal_duration", "video duration"),
}


def b_eff_backbone(tag, title):
    lat = {"3b": (8.4, 9.0, 8.0, 35.3, 17.6), "7b": (11.7, 11.7, 11.4, 46.3, 24.5), "32b": (32.2, 40.0, 31.4, 126.5, 69.6)}[tag]
    en = {"3b": (0.338, 0.362, 0.319, 1.321, 0.681), "7b": (0.529, 0.529, 0.518, 2.110, 1.104), "32b": (1.690, 2.100, 1.673, 6.717, 3.690)}[tag]
    def b(slide):
        cats = ["Direct Input", "Video CoT*", "Static Memory", "Tree d=1", "Ours"]
        ch = col_chart(slide, 1.2, 0.9, 10.9, 5.6, cats, [("Latency (s)", lat), ("Energy (kJ)", en)])
        style_series(ch.plots[0], 0, "4C72B0"); style_series(ch.plots[0], 1, "C44E52")
        add_textbox(slide, SW / 2, 6.7, title + " - measured on Jetson AGX Orin (energy in kJ on the same axis)", 13, None)
    return b

CHARTS["vt_eff_3b"] = b_eff_backbone("3b", "Qwen2.5-VL-3B")
CHARTS["vt_eff_7b"] = b_eff_backbone("7b", "Qwen2.5-VL-7B")
CHARTS["vt_eff_32b"] = b_eff_backbone("32b", "Qwen2.5-VL-32B (NF4 4-bit; bf16 OOM)")

def run():
    for name, b in CHARTS.items():
        chart_slide(name, b)
    real = matplotlib.figure.Figure.savefig
    done = set()
    def hook(self, fname, *a, **kw):
        stem = os.path.splitext(os.path.basename(str(fname)))[0]
        if str(fname).endswith(".pdf") and stem in DIAGRAMS and stem not in done:
            done.add(stem)
            try:
                convert_fig(self, stem)
            except Exception as e:
                import traceback; traceback.print_exc()
                print("CONVERT FAIL", stem, repr(e))
        return real(self, fname, *a, **kw)
    matplotlib.figure.Figure.savefig = hook
    try:
        for script in ("scripts/paper_fig_design.py", "scripts/paper_fig_wm.py"):
            path = os.path.join(R, script)
            g = {"__name__": "__main__", "__file__": path}
            exec(compile(open(path).read(), path, "exec"), g)
    finally:
        matplotlib.figure.Figure.savefig = real
    print("missing diagrams:", DIAGRAMS - done)

if __name__ == "__main__":
    run()
